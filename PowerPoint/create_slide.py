import collections.abc
from pptx import Presentation
from pptx.util import Inches # to add images 
from pptx.enum.shapes import MSO_SHAPE # to add shapes 
from pptx.enum.dml import MSO_THEME_COLOR # change shape colors 
#Graphs imports 
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_TICK_MARK
from pptx.util import Pt

pr1 = Presentation() #new presentaion 
slide1_register = pr1.slide_layouts[0] #add first slide 0 - title slide 
# 0 - title , 1 - title and content, 3 - section header 
slide2_register = pr1.slide_layouts[1]
slide3_register = pr1.slide_layouts[5]
slide4_register = pr1.slide_layouts[5]
slide5_register = pr1.slide_layouts[5]

slide1 = pr1.slides.add_slide(slide1_register)
slide2 = pr1.slides.add_slide(slide2_register)
slide3 = pr1.slides.add_slide(slide3_register)
slide4 = pr1.slides.add_slide(slide4_register)
slide5 = pr1.slides.add_slide(slide5_register)

#placeholder items in the layout 

title1 = slide1.shapes.title
subtitle1 = slide1.placeholders[1]
title2 = slide2.shapes.title 
title3 = slide3.shapes.title
title4 = slide4.shapes.title
title5 = slide5.shapes.title

#Bullets in the Slides 

bullet_point_box = slide2.shapes
bullet_point_level1 = bullet_point_box.placeholders[1]
bullet_point_level1.text = "First Bullet Point"
bullet_point_level2 = bullet_point_level1.text_frame.add_paragraph()
bullet_point_level2.text = "Sub level Second Bullet"
bullet_point_level2.level = 1 

bullet_point_level3 = bullet_point_level1.text_frame.add_paragraph()
bullet_point_level3.text = "Third Level"
bullet_point_level3.level = 2 

#pictures in the slide
img1 = "Image.PNG"

from_left = Inches(1)
from_top = Inches(2)
add_picture = slide3.shapes.add_picture(img1, from_left, from_top)

#Add shapes 
left1 = top1 = width1 = height1 = Inches(2)
add_shape1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left1, top1, width1, height1) 

left2 = Inches(6)
top2 = Inches (2)
width2 = height2 = Inches(2)
arrow1 = slide4.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left2, top2, width2, height2) 

fill_arrow1 = arrow1.fill
fill_arrow1.solid()
fill_arrow1.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_2

arrow1.rotation = 90 

#Graphs 
graph_info = CategoryChartData()
graph_info.categories = ["A", "B", "C"]
graph_info.add_series("Series 1", (15, 11, 18))

#Add the graph to the Slide 
left_graph = Inches(2)
top_graph = Inches(2)
width_graph = Inches(6)
height_graph = Inches(4)

#slide5.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left_graph,top_graph, height_graph,width_graph,graph_info)
graph1_frame = slide5.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left_graph,top_graph, height_graph,width_graph,graph_info)

graph1 = graph1_frame.chart
category_axis = graph1.category_axis
category_axis.has_major_gridlines = True
category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
category_axis.tick_labels.font.italic = True
category_axis.tick_labels.font.size = Pt(24)
#text in the slides 
title1.text = "FirstSlide"
subtitle1.text = "Just the first tutorial" 

title2.text = "Second Slide with Bullets" 
title3.text = "Picture Time !" 
title4.text = "Shapes"
title5.text = "Bar Chart"
pr1.save("First_Presentation_4.pptx") #save presenatation 